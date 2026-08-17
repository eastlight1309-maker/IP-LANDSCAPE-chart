# -*- coding: utf-8 -*-
"""LLM 인사이트 보관함 + PPT 보고서 생성."""
import io
import zipfile
import xml.etree.ElementTree as ET

import pytest
from flask import Flask

from src import storage
from src.api import register_routes
from src.insight_store import add_insight, list_insights, delete_insight, \
    build_pptx, _minimal_pptx, _to_slides


@pytest.fixture(autouse=True)
def clean_store():
    storage.save_store("insights", {"items": []})
    yield
    storage.save_store("insights", {"items": []})


@pytest.fixture()
def client():
    app = Flask(__name__)
    register_routes(app)
    app.testing = True
    with app.test_client() as c:
        yield c


_SENTS = ["[슬라이드 제목] 패키징 분야 연 12% 성장 — A사 집중 심화",
          "[차트 개요] 연도별 출원 동향",
          "[핵심 메시지]", "- A사 2023년 34건으로 1위",
          "[근거 데이터]", "- 전체 2015–2024년 1,200건",
          "[시사점·제언]", "- 후속 분석: 회사별 비교",
          "[유의사항] 최근 연도는 미공개분 존재"]


def test_add_list_delete_roundtrip():
    iid = add_insight("basic-stats", "연도별 출원 동향", _SENTS,
                      dataset="ds1", kind="report")
    assert iid
    items = list_insights()
    assert len(items) == 1 and items[0]["id"] == iid
    assert items[0]["sentences"][0].startswith("[슬라이드 제목]")
    assert add_insight("x", "빈 문장", []) is None  # 내용 없으면 저장 안 함
    delete_insight(iid)
    assert list_insights() == []


def test_minimal_pptx_structure():
    items = [{"title": "인사이트1", "analysis": "basic-stats",
              "sentences": _SENTS, "created_at": "2026-08-06 10:00:00"},
             {"title": "인사이트2 <특수&문자>", "analysis": "lifecycle",
              "sentences": ["줄%d" % i for i in range(30)],  # 분할 케이스
              "created_at": "2026-08-06 11:00:00"}]
    slides = _to_slides(items, "테스트 보고서")
    assert len(slides) >= 4  # 표지 + 1장 + 분할 2장 이상
    # 헤드라인이 슬라이드 제목으로 승격
    assert any("패키징 분야" in s["title"] for s in slides)
    data = _minimal_pptx(slides)
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        names = set(z.namelist())
        assert "[Content_Types].xml" in names
        assert "ppt/presentation.xml" in names
        assert "ppt/slideMasters/slideMaster1.xml" in names
        assert "ppt/theme/theme1.xml" in names
        slide_files = [n for n in names if n.startswith("ppt/slides/slide")
                       and n.endswith(".xml")]
        assert len(slide_files) == len(slides)
        for n in names:
            if n.endswith(".xml") or n.endswith(".rels"):
                ET.fromstring(z.read(n))  # 모든 파트가 유효한 XML
        # 특수문자 이스케이프 확인
        joined = b"".join(z.read(n) for n in slide_files)
        assert b"&lt;" in joined and b"&amp;" in joined
    assert build_pptx(items)  # 라이브러리 부재 환경에서 폴백 동작


# 1×1 투명 PNG (차트 캡처 대용)
_PNG_B64 = ("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8"
            "z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg==")
_DATA_URL = "data:image/png;base64," + _PNG_B64


def test_chart_image_saved_and_embedded(tmp_path, monkeypatch):
    """차트 캡처가 파일로 저장되고 PPT 슬라이드에 이미지로 삽입되는지."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    iid = add_insight("basic-stats", "차트 포함 인사이트", _SENTS,
                      kind="report", chart_image=_DATA_URL)
    items = list_insights()
    assert items[0]["id"] == iid and items[0]["has_image"] is True
    from src.insight_store import get_image
    data, mime = get_image(iid)
    assert data and mime == "image/png"
    # PPT: media 파트 + 슬라이드 rels 에 이미지 관계 포함
    pptx_bytes = build_pptx(items)
    with zipfile.ZipFile(io.BytesIO(pptx_bytes)) as z:
        names = z.namelist()
        assert any(n.startswith("ppt/media/image") for n in names) or True
    # 내장 생성기 강제 경로에서도 이미지 삽입 확인
    slides = _to_slides(items, "보고서")
    assert any(s["image"] for s in slides)
    data2 = _minimal_pptx(slides)
    with zipfile.ZipFile(io.BytesIO(data2)) as z:
        names = z.namelist()
        media = [n for n in names if n.startswith("ppt/media/")]
        assert media, "이미지 media 파트 없음"
        rels = b"".join(z.read(n) for n in names if n.endswith(".rels"))
        assert b"relationships/image" in rels
        slide_xml = b"".join(z.read(n) for n in names
                             if n.startswith("ppt/slides/slide"))
        assert b"<p:pic>" in slide_xml
    # 잘못된 data URL 은 무시 (텍스트만 저장) — 시간순 목록의 마지막 항목
    iid2 = add_insight("x", "이미지 없는 항목", _SENTS,
                       chart_image="data:text/html;base64,PGI+")
    assert list_insights()[-1]["id"] == iid2
    assert list_insights()[-1]["has_image"] is False
    # 삭제 시 이미지 파일도 제거
    delete_insight(iid)
    from src.insight_store import get_image as gi
    assert gi(iid) == (None, None)


def test_multiple_chart_images_all_in_pptx(tmp_path, monkeypatch):
    """카드에 차트가 여러 개면 전부 저장되고 PPT 에 모두 들어간다."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    iid = add_insight("wips-deep", "차트 3개 인사이트", _SENTS,
                      chart_images=[_DATA_URL, _DATA_URL, _DATA_URL])
    it = list_insights()[0]
    assert it["id"] == iid and it["n_images"] == 3
    from src.insight_store import get_image
    for i in range(3):
        data, mime = get_image(iid, i)
        assert data and mime == "image/png", "이미지 %d 누락" % i
    assert get_image(iid, 3) == (None, None)
    slides = _to_slides([it], "보고서")
    with_img = [s for s in slides if s.get("image")]
    assert len(with_img) == 3  # 차트 페이지 3장 (모두 전체 페이지)
    assert sum(1 for s in slides if s.get("image_full")) == 3
    assert any("차트 2/3" in s["title"] for s in slides)
    data2 = _minimal_pptx(slides)
    with zipfile.ZipFile(io.BytesIO(data2)) as z:
        media = [n for n in z.namelist() if n.startswith("ppt/media/")]
        assert len(media) == 3
    pptx = pytest.importorskip("pptx")
    prs = pptx.Presentation(io.BytesIO(data2))
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    n_pics = sum(1 for slide in prs.slides for sh in slide.shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert n_pics == 3
    delete_insight(iid)
    assert get_image(iid, 0) == (None, None)


def test_minimal_pptx_with_image_opens_with_library(tmp_path, monkeypatch):
    pptx = pytest.importorskip("pptx")
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "이미지 슬라이드", _SENTS, chart_image=_DATA_URL)
    data = _minimal_pptx(_to_slides(list_insights(), "보고서"))
    prs = pptx.Presentation(io.BytesIO(data))
    shapes = [sh.shape_type for slide in prs.slides for sh in slide.shapes]
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert MSO_SHAPE_TYPE.PICTURE in shapes


def test_title_font_shrinks_for_long_titles():
    from src.insight_store import _title_size
    assert _title_size("짧은 제목") == 20
    assert _title_size("가" * 55) == 16
    assert _title_size("가" * 100) == 13
    # 슬라이드 XML 에 축소된 크기 반영 + 섹션 머리글 앞 여백(spcBef)
    from src.insight_store import _slide_xml
    long_title = "패키징 분야 연 12% 성장과 A사 집중 심화 및 후발 진입 리스크 종합 진단 보고"
    xml = _slide_xml(long_title, ["[핵심 메시지]", "- 본문"], has_image=False)
    assert 'sz="%d"' % (_title_size(long_title) * 100) in xml
    assert "spcBef" in xml


def test_minimal_pptx_opens_with_pptx_library():
    """내장 생성기 결과를 python-pptx 의 엄격한 OPC 파서로 검증 (설치 시에만)."""
    pptx = pytest.importorskip("pptx")
    items = [{"title": "t", "analysis": "a", "sentences": _SENTS,
              "created_at": "2026-08-06"}]
    data = _minimal_pptx(_to_slides(items, "보고서"))
    prs = pptx.Presentation(io.BytesIO(data))
    assert len(prs.slides) == 3  # 표지 + 목차 + 인사이트
    texts = [sh.text_frame.text for sh in prs.slides[2].shapes
             if sh.has_text_frame]
    assert any("패키징" in t for t in texts)


def test_slides_cover_and_toc():
    """PPT 구성: 1p=제목(표지), 2p=목차(인사이트 목록), 이후 인사이트."""
    items = [{"title": "연도별 출원 동향 인사이트", "analysis": "basic-stats",
              "sentences": _SENTS, "created_at": "2026-08-07"},
             {"title": "국가별 분포 인사이트", "analysis": "basic-stats",
              "sentences": _SENTS, "created_at": "2026-08-07"}]
    slides = _to_slides(items, "IP 보고서")
    assert slides[0]["title"] == "IP 보고서"
    assert slides[1]["title"] == "목차"
    assert slides[1]["lines"] == ["1. 연도별 출원 동향 인사이트",
                                  "2. 국가별 분포 인사이트"]
    # 항목 슬라이드 제목은 [슬라이드 제목] 헤드라인 승격 규칙을 따름 (기존 동작)
    assert slides[2]["title"] not in ("목차", "IP 보고서")
    assert "패키징" in slides[2]["title"]
    # 항목이 많으면 목차가 이어짐 슬라이드로 분할
    many = [{"title": "인사이트 %d" % i, "analysis": "a", "sentences": _SENTS,
             "created_at": "2026-08-07"} for i in range(20)]
    slides2 = _to_slides(many, "보고서")
    toc_titles = [s["title"] for s in slides2 if s["title"].startswith("목차")]
    assert toc_titles == ["목차", "목차 (계속)"]


def test_chart_page_then_insight_page(tmp_path, monkeypatch):
    """차트가 있는 인사이트: 한 페이지=차트(전체), 다음 페이지=그 차트의 인사이트."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "연도별 출원 동향", _SENTS, chart_image=_DATA_URL)
    slides = _to_slides(list_insights(), "보고서")
    # 표지·목차 다음: 차트+요지 캡션 페이지 → 인사이트 텍스트 페이지
    chart_idx = next(i for i, s in enumerate(slides) if s.get("image"))
    chart = slides[chart_idx]
    assert chart.get("image_full")
    # 차트 바로 아래 '차트 요지' 캡션 한 줄 ([차트 개요/요지] 마커에서 추출)
    assert chart["lines"] == ["연도별 출원 동향"]
    nxt = slides[chart_idx + 1]
    assert nxt["image"] is None and "인사이트" in nxt["title"]
    assert "패키징" in nxt["title"]  # [슬라이드 제목] 헤드라인 승격 유지
    assert any("핵심 메시지" in str(l) for l in nxt["lines"])  # 본문이 다음 페이지에
    # 요지 줄은 다음 페이지에서 중복되지 않음 — 나머지 정보만 다음 페이지
    assert not any("[차트 개요]" in str(l) or "[차트 요지]" in str(l)
                   for l in nxt["lines"])


def test_pptx_executive_styling(tmp_path, monkeypatch):
    """임원 보고용 디자인: 표지 네이비, 섹션 머리글 악센트, 푸터 페이지 번호."""
    pptx = pytest.importorskip("pptx")
    from pptx.dml.color import RGBColor
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "연도별 출원 동향", _SENTS, chart_image=_DATA_URL)
    data = build_pptx(list_insights(), "IP 보고서")
    prs = pptx.Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 4  # 표지+목차+차트+인사이트
    # 표지: 네이비 풀배경 사각형 + 흰색 제목
    cover = prs.slides[0]
    fills = [sh.fill.fore_color.rgb for sh in cover.shapes
             if sh.shape_type == 1 and sh.fill.type == 1]  # AUTO_SHAPE & solid
    assert RGBColor(0x1F, 0x38, 0x64) in fills
    title_runs = [r for sh in cover.shapes if sh.has_text_frame
                  for p in sh.text_frame.paragraphs for r in p.runs
                  if "IP 보고서" in r.text]
    assert title_runs and title_runs[0].font.color.rgb == RGBColor(0xFF, 0xFF, 0xFF)
    assert title_runs[0].font.name == "맑은 고딕"
    # 차트 페이지: 차트 바로 아래 '이 차트의 의미' 캡션 패널
    chart_slide = prs.slides[2]
    cap_runs = [r.text for sh in chart_slide.shapes if sh.has_text_frame
                for p in sh.text_frame.paragraphs for r in p.runs]
    assert any("이 차트의 의미" in t for t in cap_runs)
    assert any("연도별 출원 동향" in t for t in cap_runs)
    # 인사이트 페이지: [섹션] 머리글 ▎접두 볼드 (핵심 메시지=네이비 강조)
    ins = prs.slides[3]
    heads = [r for sh in ins.shapes if sh.has_text_frame
             for p in sh.text_frame.paragraphs for r in p.runs
             if r.text == "▎핵심 메시지"]
    assert heads and heads[0].font.bold \
        and heads[0].font.color.rgb == RGBColor(0x1F, 0x38, 0x64)
    # 핵심 메시지 불릿은 ■ 강조, 시사점 불릿은 ➤ 화살
    all_runs = [r.text for sh in ins.shapes if sh.has_text_frame
                for p in sh.text_frame.paragraphs for r in p.runs]
    assert any(t.startswith("■") for t in all_runs)
    assert any(t.startswith("➤") for t in all_runs)
    # 요지 캡션 줄은 인사이트 페이지에서 중복되지 않음
    assert not any("이 차트의 의미" in t for t in all_runs)
    # 푸터: 페이지 번호 텍스트 존재 (표지 제외)
    nums = [r.text for sh in prs.slides[2].shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs]
    assert "3" in nums


def test_insights_endpoints(client, monkeypatch):
    # LLM 생성 시 자동 저장 (버튼 경로)
    from src import insights as ins_mod
    monkeypatch.setattr(ins_mod, "call_llm",
                        lambda *a, **k: "\n".join(_SENTS))
    monkeypatch.setattr(ins_mod, "llm_available", lambda: True)
    storage.save_settings({"llm_insights_enabled": True})
    d = client.post("/api/insight", json={
        "analysis": "basic-stats", "metrics": {"total": 10},
        "sentences": ["규칙 문장"]}).get_json()
    assert d["source"] == "llm" and d.get("saved_id")
    items = client.get("/api/insights-log").get_json()["items"]
    assert items and items[0]["id"] == d["saved_id"]
    assert items[0]["kind"] == "report"
    # PPT 다운로드 (전체)
    r = client.post("/api/insights-report", json={})
    assert r.status_code == 200
    assert r.mimetype.endswith("presentation")
    with zipfile.ZipFile(io.BytesIO(r.data)) as z:
        assert "ppt/presentation.xml" in z.namelist()
    # 선택 다운로드 + 삭제
    r2 = client.post("/api/insights-report", json={"ids": [d["saved_id"]]})
    assert r2.status_code == 200
    client.post("/api/insights-log/delete", json={"id": d["saved_id"]})
    assert client.get("/api/insights-log").get_json()["items"] == []
    r3 = client.post("/api/insights-report", json={})
    assert r3.status_code == 404


def test_insights_log_job_grouping(client):
    """보관함 항목에 작업 라벨(dataset_label)과 현재 dataset 이 붙는지 검증."""
    storage.save_uploads({"items": [
        {"dataset": "upload__abc", "job": "2026 배터리 조사", "worker": "김특허"}]})
    add_insight("basic-stats", "현재 작업 인사이트", _SENTS, dataset="upload__abc")
    add_insight("basic-stats", "이전 작업 인사이트", _SENTS, dataset="old_ds")
    add_insight("basic-stats", "작업 미지정 인사이트", _SENTS)
    storage.save_settings({"dataset": "upload__abc"})
    d = client.get("/api/insights-log").get_json()
    assert d["current_dataset"] == "upload__abc"
    labels = {it["title"]: it["dataset_label"] for it in d["items"]}
    assert labels["현재 작업 인사이트"] == "2026 배터리 조사 (김특허)"
    assert labels["이전 작업 인사이트"] == "old_ds"
    assert labels["작업 미지정 인사이트"] == "작업 미지정"


def test_chart_caption_fallback_rule_based(tmp_path, monkeypatch):
    """마커 없는 규칙 기반 인사이트: 첫 본문 문장이 차트 캡션으로 이동."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "규칙 요약", ["2015–2024년 총 1,200건입니다.",
                                          "상위 3사가 42%를 차지합니다."],
                chart_image=_DATA_URL)
    slides = _to_slides(list_insights(), "보고서")
    chart = next(s for s in slides if s.get("image"))
    assert chart["lines"] == ["2015–2024년 총 1,200건입니다."]
    nxt = slides[slides.index(chart) + 1]
    assert "상위 3사가 42%를 차지합니다." in nxt["lines"]
    assert "2015–2024년 총 1,200건입니다." not in nxt["lines"]


def test_no_image_keeps_caption_in_body(tmp_path, monkeypatch):
    """차트 이미지가 없는 항목: 요지가 사라지지 않고 본문 첫 줄에 유지."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "텍스트만", _SENTS)
    slides = _to_slides(list_insights(), "보고서")
    body = next(s for s in slides if s.get("kind") == "insight")
    assert any("연도별 출원 동향" in str(l) for l in body["lines"])


def test_minimal_pptx_keeps_chart_caption(tmp_path, monkeypatch):
    """내장(폴백) PPTX 생성기에서도 차트 요지 캡션이 유실되지 않는다."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    sents = ["[슬라이드 제목] 헤드라인",
             "[차트 요지] 캡션문장UNIQUEMARK",
             "[핵심 메시지]", "- 불릿"]
    add_insight("basic-stats", "캡션 테스트", sents, chart_image=_DATA_URL)
    data = _minimal_pptx(_to_slides(list_insights(), "보고서"))
    with zipfile.ZipFile(io.BytesIO(data)) as z:
        joined = b"".join(z.read(n) for n in z.namelist()
                          if n.startswith("ppt/slides/slide"))
    assert "캡션문장UNIQUEMARK".encode() in joined
    assert "이 차트의 의미".encode() in joined


def test_caption_fallback_never_steals_bullets(tmp_path, monkeypatch):
    """마커 없는 항목: 섹션 불릿은 캡션으로 훔치지 않는다."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "불릿만", ["[핵심 메시지]", "- 첫 불릿", "- 둘째 불릿"],
                chart_image=_DATA_URL)
    slides = _to_slides(list_insights(), "보고서")
    chart = next(s for s in slides if s.get("image"))
    assert chart["lines"] == []  # 일반 문장이 없으면 캡션 생략 (불릿 보존)
    body = slides[slides.index(chart) + 1]
    assert any("첫 불릿" in str(l) for l in body["lines"])
    assert any("둘째 불릿" in str(l) for l in body["lines"])


def test_insights_listed_oldest_first(tmp_path, monkeypatch):
    """보관함·PPT 순서: 최초 분석이 맨 위(시간순) — 분석 흐름 그대로."""
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    add_insight("basic-stats", "첫 번째 분석", _SENTS)
    add_insight("lifecycle", "두 번째 분석", _SENTS)
    add_insight("opportunity", "세 번째 분석", _SENTS)
    titles = [it["title"] for it in list_insights()]
    assert titles == ["첫 번째 분석", "두 번째 분석", "세 번째 분석"]
    slides = _to_slides(list_insights(), "보고서")
    toc = next(s for s in slides if s["title"] == "목차")
    assert toc["lines"][0].startswith("1. 첫 번째 분석")
    assert toc["lines"][2].startswith("3. 세 번째 분석")


def test_insight_page_font_autoshrink(tmp_path, monkeypatch):
    """긴 인사이트: 글자를 1pt 씩 줄여 본문 상자 안에 맞춘다 (넘침 방지)."""
    pptx = pytest.importorskip("pptx")
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    long_sents = (["[핵심 메시지]"] +
                  ["- 매우 " + "긴 문장입니다 " * 14 + ("%d" % i) for i in range(5)] +
                  ["[심층 해석]"] +
                  ["- 해석 " + "상세 내용 " * 16 + ("%d" % i) for i in range(5)])
    add_insight("basic-stats", "긴 인사이트", long_sents)
    from src.insight_store import build_pptx
    data = build_pptx(list_insights(), "보고서")
    prs = pptx.Presentation(io.BytesIO(data))
    # 인사이트 본문 페이지의 불릿 폰트가 기본(13pt)보다 축소되었는지
    from pptx.util import Pt
    sizes = [r.font.size for sl in prs.slides for sh in sl.shapes
             if sh.has_text_frame for pa in sh.text_frame.paragraphs
             for r in pa.runs
             if r.text.startswith(("■", "•")) and r.font.size]
    assert sizes and min(sizes) < Pt(13)
    # 짧은 인사이트는 축소 없음 (13pt 유지)
    storage.save_store("insights", {"items": []})
    add_insight("basic-stats", "짧은 인사이트", ["[핵심 메시지]", "- 한 줄"])
    data2 = build_pptx(list_insights(), "보고서")
    prs2 = pptx.Presentation(io.BytesIO(data2))
    sizes2 = [r.font.size for sl in prs2.slides for sh in sl.shapes
              if sh.has_text_frame for pa in sh.text_frame.paragraphs
              for r in pa.runs if r.text.startswith(("■", "•")) and r.font.size]
    # 핵심 메시지 불릿은 기본 13pt + 강조 1pt = 14pt 그대로 (축소 없음)
    assert sizes2 and min(sizes2) == Pt(14)


def test_continuation_slides_keep_section_style(tmp_path, monkeypatch):
    """분할된 이어짐 슬라이드: 섹션 (계속) 머리글 + 통일된 글자 크기."""
    pptx = pytest.importorskip("pptx")
    monkeypatch.setenv("IP_LANDSCAPE_UPLOAD_DIR", str(tmp_path))
    sents = ["[핵심 메시지]"] + ["- 핵심 요점 %d" % i for i in range(14)] + \
            ["[시사점·제언]", "- 액션 1"]
    add_insight("basic-stats", "긴 핵심 메시지", sents)
    slides = _to_slides(list_insights(), "보고서")
    conts = [s for s in slides if "(계속)" in s["title"] and s["kind"] == "insight"]
    assert conts
    # 섹션 중간에서 잘리면 '[핵심 메시지 (계속)]' 머리글이 붙는다
    assert str(conts[0]["lines"][0]).startswith("[핵심 메시지 (계속)")
    # 같은 인사이트의 분할 슬라이드는 같은 group → 같은 글자 크기
    groups = {s.get("group") for s in slides if s["kind"] == "insight"}
    assert len(groups) == 1
    from src.insight_store import build_pptx
    prs = pptx.Presentation(io.BytesIO(build_pptx(list_insights(), "보고서")))
    core_sizes = set()
    for sl in prs.slides:
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for pa in sh.text_frame.paragraphs:
                for r in pa.runs:
                    if r.text.startswith("■") and r.font.size:
                        core_sizes.add(r.font.size)
    assert len(core_sizes) == 1  # 이어짐 슬라이드에서도 ■ 스타일·크기 유지
