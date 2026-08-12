# -*- coding: utf-8 -*-
"""
insight_store.py — LLM 인사이트 보관함 + PPT 보고서 생성.

보관함:
- LLM 이 생성한 인사이트(버튼 결과·챗 답변)를 storage("insights")에 자동 저장한다
  (Dataiku 프로젝트 변수 / 로컬 JSON — Backend 재시작 후에도 유지, 최근 300건).
- 항목: {id, kind(report|chat), analysis, title, question?, sentences[], dataset,
        created_at}

PPT 보고서:
- 저장된 인사이트를 .pptx 로 내보낸다. python-pptx 가 설치되어 있으면 사용하고,
  없으면 외부 의존성 없는 내장 OOXML 생성기(_minimal_pptx)로 생성한다 —
  표지 1장 + 인사이트당 1장(길면 이어짐 슬라이드), 텍스트 전용 16:9.
"""
import base64
import io
import logging
import os
import re
import time
import uuid
import zipfile
from xml.sax.saxutils import escape

from src import storage

logger = logging.getLogger("ip_landscape")

_MAX_ITEMS = 300
_LINES_PER_SLIDE = 13
_MAX_IMAGE_MB = 4
_DATAURL_RE = re.compile(r"^data:image/(png|jpeg);base64,(.+)$", re.DOTALL)


_MAX_IMAGES = 6  # 카드 하나에서 캡처·저장하는 차트 수 상한


def _save_chart_image(insight_id, chart_image, idx=0):
    """프론트가 보낸 차트 캡처(data URL) → PNG/JPEG 파일 저장. 파일명 또는 None."""
    m = _DATAURL_RE.match(str(chart_image or "").strip())
    if not m:
        return None
    try:
        raw = base64.b64decode(m.group(2), validate=False)
    except Exception:
        return None
    if not raw or len(raw) > _MAX_IMAGE_MB * 1024 * 1024:
        return None
    ext = "png" if m.group(1) == "png" else "jpg"
    fname = ("%s.%s" % (insight_id, ext) if idx == 0
             else "%s_%d.%s" % (insight_id, idx, ext))
    try:
        with open(os.path.join(storage.insight_image_dir(), fname), "wb") as fh:
            fh.write(raw)
        return fname
    except OSError as e:
        logger.warning("인사이트 이미지 저장 실패: %s", e)
        return None


def _image_paths(entry):
    """항목의 차트 이미지 파일 경로 목록 (존재하는 것만, 저장 순서 유지)."""
    fnames = entry.get("image_files")
    if not fnames:
        fnames = [entry.get("image_file")] if entry.get("image_file") else []
    out = []
    for fname in fnames:
        path = os.path.join(storage.insight_image_dir(), str(fname))
        if os.path.exists(path):
            out.append(path)
    return out


def _image_path(entry):
    paths = _image_paths(entry)
    return paths[0] if paths else None


def _remove_image(entry):
    for path in _image_paths(entry):
        try:
            os.remove(path)
        except OSError:
            pass


def add_insight(analysis, title, sentences, dataset=None, kind="report",
                question=None, chart_image=None, chart_images=None,
                owner=None):
    """LLM 인사이트 저장 (+차트 캡처 이미지들 — PPT 삽입용). 반환: 항목 id.

    chart_images: 카드의 모든 차트 캡처(data URL 목록) — PPT 에 전부 들어간다.
    chart_image: 구버전 단일 캡처 (chart_images 미지정 시 사용).
    """
    sentences = [str(s).strip() for s in (sentences or []) if str(s).strip()][:40]
    if not sentences:
        return None
    uid = uuid.uuid4().hex[:10]
    images_in = [img for img in (chart_images or []) if img] or \
        ([chart_image] if chart_image else [])
    image_files = []
    for i, img in enumerate(images_in[:_MAX_IMAGES]):
        fname = _save_chart_image(uid, img, idx=len(image_files))
        if fname:
            image_files.append(fname)
    entry = {
        "id": uid, "kind": kind,
        "analysis": str(analysis or "")[:60],
        "title": str(title or analysis or "인사이트")[:160],
        "question": (str(question)[:200] if question else None),
        "sentences": sentences,
        "dataset": (str(dataset)[:80] if dataset else None),
        "created_at": time.strftime("%Y-%m-%d %H:%M:%S"),
        "owner": (str(owner).strip()[:60] if owner else None),
        "image_file": image_files[0] if image_files else None,
        "image_files": image_files,
    }
    data = storage.load_store("insights")
    items = list(data.get("items") or [])
    items.insert(0, entry)
    for evicted in items[_MAX_ITEMS:]:  # 상한 초과분의 이미지 파일 정리
        _remove_image(evicted)
    storage.save_store("insights", {"items": items[:_MAX_ITEMS]})
    return entry["id"]


def list_insights():
    items = list(storage.load_store("insights").get("items") or [])
    for it in items:
        paths = _image_paths(it)
        it["has_image"] = bool(paths)
        it["n_images"] = len(paths)
    return items


def get_image(insight_id, idx=0):
    """항목의 idx 번째 차트 이미지 (bytes, mimetype) 또는 (None, None)."""
    for it in storage.load_store("insights").get("items") or []:
        if str(it.get("id")) == str(insight_id):
            paths = _image_paths(it)
            try:
                idx = int(idx or 0)
            except (TypeError, ValueError):
                idx = 0
            if 0 <= idx < len(paths):
                path = paths[idx]
                with open(path, "rb") as fh:
                    return fh.read(), ("image/png" if path.endswith(".png")
                                       else "image/jpeg")
    return None, None


def delete_insight(insight_id):
    data = storage.load_store("insights")
    items = list(data.get("items") or [])
    for it in items:
        if str(it.get("id")) == str(insight_id):
            _remove_image(it)
    items = [it for it in items if str(it.get("id")) != str(insight_id)]
    storage.save_store("insights", {"items": items})
    return True


def get_insights(ids=None):
    items = list_insights()
    if not ids:
        return items
    wanted = set(map(str, ids))
    return [it for it in items if str(it.get("id")) in wanted]


# ---------------------------------------------------------------------------
# PPTX 생성
# ---------------------------------------------------------------------------
def build_pptx(items, report_title="IP Landscape 인사이트 보고서"):
    """인사이트 목록 → .pptx 바이트. python-pptx 우선, 내장 생성기 폴백."""
    slides = _to_slides(items, report_title)
    try:
        return _pptx_via_library(slides)
    except ImportError:
        return _minimal_pptx(slides)


def _to_slides(items, report_title):
    """항목 → [{"title","lines","image","ext"}]. 긴 항목은 이어짐 슬라이드로 분할.

    구성: ① 표지(제목) ② 목차(인사이트 목록) ③ 인사이트마다
    [차트 전체 페이지] → [다음 페이지에 그 차트의 인사이트 텍스트] 순서.
    카드에 차트가 여러 개인 항목(구버전 카드 단위 캡처)은 나머지 차트도
    "차트 k/n" 전체 페이지로 이어서 들어간다.
    """
    slides = [{"title": report_title, "image": None, "ext": None, "kind": "cover",
               "lines": ["생성일: %s" % time.strftime("%Y-%m-%d"),
                         "포함 인사이트: %d건" % len(items),
                         "", "본 보고서의 지표는 특허 데이터 기반 통계 신호이며 "
                         "법률 자문(FTO·유효성 판단)을 대체하지 않습니다."]}]
    # ② 목차 — 인사이트 제목 목록 (많으면 이어짐 슬라이드로 분할)
    toc_lines = ["%d. %s" % (i + 1,
                             str(it.get("title") or it.get("analysis") or "인사이트")[:80])
                 for i, it in enumerate(items)]
    for start in range(0, len(toc_lines), _LINES_PER_SLIDE):
        slides.append({"title": "목차" if start == 0 else "목차 (계속)",
                       "lines": toc_lines[start:start + _LINES_PER_SLIDE],
                       "image": None, "ext": None, "kind": "toc"})
    for it in items:
        title = str(it.get("title") or it.get("analysis") or "인사이트")
        # 첫 줄이 [슬라이드 제목] 헤드라인이면 그 내용을 슬라이드 제목으로 사용
        lines = list(it.get("sentences") or [])
        if lines and lines[0].startswith("[슬라이드 제목]"):
            title = lines[0].replace("[슬라이드 제목]", "").strip() or title
            lines = lines[1:]
        meta_line = "· %s · %s%s" % (it.get("analysis", ""),
                                     it.get("created_at", ""),
                                     (" · Q: %s" % it["question"])
                                     if it.get("question") else "")
        lines = [meta_line] + lines
        images = []
        for path in _image_paths(it):
            try:
                with open(path, "rb") as fh:
                    images.append((fh.read(),
                                   "png" if path.endswith(".png") else "jpg"))
            except OSError:
                continue
        # ① 차트 페이지: 첫 차트를 한 페이지 가득 배치
        if images:
            img0, ext0 = images[0]
            slides.append({"title": title[:120], "lines": [],
                           "image": img0, "ext": ext0, "image_full": True,
                           "kind": "chart"})
        # ② 다음 페이지: 그 차트의 인사이트 텍스트 (길면 이어짐 분할)
        for start in range(0, len(lines), _LINES_PER_SLIDE):
            chunk = lines[start:start + _LINES_PER_SLIDE]
            if start == 0:
                t = (title + " — 인사이트") if images else title
            else:
                t = title[:60] + (" — 인사이트 (계속)" if images else " (계속)")
            slides.append({"title": t[:120], "lines": chunk,
                           "image": None, "ext": None, "kind": "insight"})
        # ③ 카드에 차트가 여러 개인 항목: 나머지 차트도 전체 페이지로 포함
        for k, (img, ext) in enumerate(images[1:], start=2):
            slides.append({"title": ("%s — 차트 %d/%d" % (title[:100], k,
                                                        len(images)))[:120],
                           "lines": [], "image": img, "ext": ext,
                           "image_full": True, "kind": "chart"})
    return slides


def _pptx_via_library(slides):
    """python-pptx 기반 임원 보고용 디자인.

    구성 원칙: 표지=네이비 풀배경, 본문=흰 배경 + 상단 타이틀 밴드(악센트 룰),
    차트 페이지=이미지 비율 유지 중앙 배치, 인사이트 페이지=[섹션] 머리글
    악센트 컬러 + 불릿 들여쓰기, 전 페이지 하단 푸터(보고서명 · 페이지 번호).
    """
    from pptx import Presentation  # noqa — 미설치 시 ImportError → 내장 생성기
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    NAVY = RGBColor(0x1F, 0x38, 0x64)     # 표지·제목
    ACCENT = RGBColor(0x2E, 0x74, 0xB5)   # 악센트(룰·섹션 머리글)
    TEXT = RGBColor(0x33, 0x3F, 0x4E)     # 본문
    SOFT = RGBColor(0x8A, 0x99, 0xA8)     # 메타·푸터
    COVER_SUB = RGBColor(0xC9, 0xD7, 0xEA)
    PAGE_W, PAGE_H = Inches(13.333), Inches(7.5)
    KOR_FONT = "맑은 고딕"

    def _font(run, size, bold=False, color=TEXT, name=KOR_FONT):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = name
        # 한글은 eastAsia typeface 를 별도 지정해야 적용된다
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn("a:ea"))
        if ea is None:
            ea = rpr.makeelement(qn("a:ea"), {})
            rpr.append(ea)
        ea.set("typeface", name)

    def _para(tf, first_used):
        return tf.paragraphs[0] if not first_used[0] else tf.add_paragraph()

    def _rect(slide, x, y, w, h, color):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
        shp.line.fill.background()
        shp.shadow.inherit = False
        return shp

    def _footer(slide, report_title, page_no):
        _rect(slide, Inches(0.55), Inches(7.12), Inches(12.23), Emu(9525), SOFT)
        fb = slide.shapes.add_textbox(Inches(0.55), Inches(7.14),
                                      Inches(10.0), Inches(0.32))
        p = fb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = str(report_title)[:70]
        _font(r, 9, color=SOFT)
        nb = slide.shapes.add_textbox(Inches(12.0), Inches(7.14),
                                      Inches(0.8), Inches(0.32))
        p2 = nb.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        r2 = p2.add_run()
        r2.text = str(page_no)
        _font(r2, 9, color=SOFT)

    def _title_band(slide, title):
        # 좌측 악센트 블록 + 제목 + 하단 얇은 룰
        _rect(slide, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.62), ACCENT)
        tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.32),
                                      Inches(11.9), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = str(title)
        _font(r, _title_size(title), bold=True, color=NAVY)
        _rect(slide, Inches(0.55), Inches(1.18), Inches(12.23), Emu(12700), ACCENT)

    def _add_picture_fit(slide, img_bytes, box_x, box_y, box_w, box_h):
        """이미지를 비율 유지로 상자 안에 최대 크기 배치 (왜곡 방지)."""
        pic = slide.shapes.add_picture(io.BytesIO(img_bytes), box_x, box_y,
                                       width=box_w)  # 폭 기준 → 높이 자동
        if pic.height > box_h:  # 세로가 넘치면 높이 기준으로 재조정
            ratio = box_h / float(pic.height)
            pic.width = int(pic.width * ratio)
            pic.height = int(box_h)
        pic.left = int(box_x + (box_w - pic.width) / 2)
        pic.top = int(box_y + (box_h - pic.height) / 2)
        return pic

    prs = Presentation()
    prs.slide_width, prs.slide_height = PAGE_W, PAGE_H
    blank = prs.slide_layouts[6]
    report_title = slides[0]["title"] if slides else "IP Landscape 보고서"

    for idx, sl in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        kind = sl.get("kind") or ("cover" if idx == 0 else "insight")

        if kind == "cover":
            _rect(slide, 0, 0, PAGE_W, PAGE_H, NAVY)
            _rect(slide, Inches(0.9), Inches(2.55), Inches(1.6), Emu(38100), ACCENT)
            tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.75),
                                          Inches(11.5), Inches(1.8))
            tb.text_frame.word_wrap = True
            r = tb.text_frame.paragraphs[0].add_run()
            r.text = str(sl["title"])
            _font(r, 36, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
            sb = slide.shapes.add_textbox(Inches(0.92), Inches(4.6),
                                          Inches(11.5), Inches(2.2))
            stf = sb.text_frame
            stf.word_wrap = True
            used = [False]
            for line in sl["lines"]:
                s = str(line)
                if not s:
                    continue
                p = _para(stf, used)
                used[0] = True
                r = p.add_run()
                r.text = s
                small = s.startswith("본 보고서")
                _font(r, 11 if small else 15,
                      color=SOFT if small else COVER_SUB)
                p.space_after = Pt(6)
            continue

        _title_band(slide, sl["title"])
        has_img = bool(sl.get("image"))

        if has_img and sl.get("image_full"):
            # 차트 페이지: 비율 유지 최대 배치
            _add_picture_fit(slide, sl["image"], Inches(0.7), Inches(1.4),
                             Inches(11.93), Inches(5.55))
            _footer(slide, report_title, idx + 1)
            continue
        if has_img:  # (구버전 좌图우문 레이아웃 항목 호환)
            _add_picture_fit(slide, sl["image"], Inches(0.55), Inches(1.5),
                             Inches(6.9), Inches(5.3))
            body_x, body_w, fsize = Inches(7.7), Inches(5.1), 12
        else:
            body_x, body_w, fsize = Inches(0.85), Inches(11.9), 13

        body = slide.shapes.add_textbox(body_x, Inches(1.5), body_w, Inches(5.45))
        tf = body.text_frame
        tf.word_wrap = True
        used = [False]
        for i, line in enumerate(sl["lines"]):
            s = str(line)
            p = _para(tf, used)
            used[0] = True
            r = p.add_run()
            if kind == "toc":
                num, _, rest = s.partition(". ")
                if rest:
                    r.text = num + "."
                    _font(r, 14, bold=True, color=ACCENT)
                    r2 = p.add_run()
                    r2.text = "  " + rest
                    _font(r2, 14, color=TEXT)
                else:
                    r.text = s
                    _font(r, 14, color=TEXT)
                p.space_after = Pt(9)
                continue
            if i == 0 and s.startswith("·"):
                r.text = s  # 메타 줄 (분석명·생성일)
                _font(r, 9, color=SOFT)
                p.space_after = Pt(6)
            elif s.startswith("["):
                r.text = s.strip("[]").strip() if s.endswith("]") else s
                _font(r, fsize + 2, bold=True, color=ACCENT)
                p.space_before = Pt(12)
                p.space_after = Pt(3)
            elif s.startswith(("-", "·", "•")):
                r.text = "•  " + s.lstrip("-·• ").strip()
                _font(r, fsize, color=TEXT)
                p.level = 1
                p.space_after = Pt(4)
            else:
                r.text = s
                _font(r, fsize, color=TEXT)
                p.space_after = Pt(4)
        _footer(slide, report_title, idx + 1)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---- 내장 OOXML 생성기 (외부 의존성 없음, 텍스트 전용 16:9) ----------------
_CT = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Default Extension="png" ContentType="image/png"/>
<Default Extension="jpg" ContentType="image/jpeg"/>
<Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
<Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
<Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
<Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
%s</Types>"""

_ROOT_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
</Relationships>"""

_NS = ('xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
       'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" '
       'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"')

_THEME = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="T">
<a:themeElements><a:clrScheme name="C"><a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1>
<a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="1F3B54"/></a:dk2>
<a:lt2><a:srgbClr val="EEF2F6"/></a:lt2><a:accent1><a:srgbClr val="4E79A7"/></a:accent1>
<a:accent2><a:srgbClr val="F28E2B"/></a:accent2><a:accent3><a:srgbClr val="59A14F"/></a:accent3>
<a:accent4><a:srgbClr val="E15759"/></a:accent4><a:accent5><a:srgbClr val="76B7B2"/></a:accent5>
<a:accent6><a:srgbClr val="EDC948"/></a:accent6><a:hlink><a:srgbClr val="1668A8"/></a:hlink>
<a:folHlink><a:srgbClr val="800080"/></a:folHlink></a:clrScheme>
<a:fontScheme name="F"><a:majorFont><a:latin typeface="Malgun Gothic"/><a:ea typeface="Malgun Gothic"/><a:cs typeface=""/></a:majorFont>
<a:minorFont><a:latin typeface="Malgun Gothic"/><a:ea typeface="Malgun Gothic"/><a:cs typeface=""/></a:minorFont></a:fontScheme>
<a:fmtScheme name="S"><a:fillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
<a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:fillStyleLst>
<a:lnStyleLst><a:ln w="6350"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
<a:ln w="12700"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln>
<a:ln w="19050"><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:ln></a:lnStyleLst>
<a:effectStyleLst><a:effectStyle><a:effectLst/></a:effectStyle><a:effectStyle><a:effectLst/></a:effectStyle>
<a:effectStyle><a:effectLst/></a:effectStyle></a:effectStyleLst>
<a:bgFillStyleLst><a:solidFill><a:schemeClr val="phClr"/></a:solidFill>
<a:solidFill><a:schemeClr val="phClr"/></a:solidFill><a:solidFill><a:schemeClr val="phClr"/></a:solidFill></a:bgFillStyleLst>
</a:fmtScheme></a:themeElements></a:theme>"""

_MASTER = ("""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster %s><p:cSld><p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr/></p:spTree></p:cSld>
<p:clrMap bg1="lt1" tx1="dk1" bg2="lt2" tx2="dk2" accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" hlink="hlink" folHlink="folHlink"/>
<p:sldLayoutIdLst><p:sldLayoutId id="2147483649" r:id="rId1"/></p:sldLayoutIdLst>
</p:sldMaster>""" % _NS)

_MASTER_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>"""

_LAYOUT = ("""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout %s type="blank"><p:cSld><p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr/></p:spTree></p:cSld>
<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>""" % _NS)

_LAYOUT_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="../slideMasters/slideMaster1.xml"/>
</Relationships>"""

_SLIDE_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
%s</Relationships>"""

_IMG_REL = ('<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/'
            'officeDocument/2006/relationships/image" Target="../media/%s"/>')


def _picture_xml(shape_id, x, y, w, h):
    """차트 캡처 이미지 pic 요소 (r:embed=rId2)."""
    return ('<p:pic><p:nvPicPr><p:cNvPr id="%d" name="chart"/>'
            '<p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
            '<p:blipFill><a:blip r:embed="rId2"/><a:stretch><a:fillRect/>'
            '</a:stretch></p:blipFill>'
            '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
            '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
            % (shape_id, x, y, w, h))


def _title_size(title):
    """제목 길이에 따라 글자 크기 축소 (화면 밖으로 나가지 않게)."""
    n = len(str(title))
    if n <= 40:
        return 20
    if n <= 70:
        return 16
    return 13


def _textbox(shape_id, name, x, y, w, h, paragraphs):
    """EMU 좌표 텍스트박스 sp XML.

    paragraphs: [(text, size_pt, bold[, space_before_pt])] — 섹션 머리글 앞에
    여백을 줘 그래프·본문과 조화롭게 읽히도록 한다.
    """
    paras = []
    for para in paragraphs:
        text, size, bold = para[0], para[1], para[2]
        spc = para[3] if len(para) > 3 else 0
        t = escape(str(text)) or " "
        ppr = ('<a:pPr><a:spcBef><a:spcPts val="%d"/></a:spcBef></a:pPr>'
               % int(spc * 100)) if spc else "<a:pPr/>"
        paras.append(
            '<a:p>%s<a:r><a:rPr lang="ko-KR" sz="%d" b="%d" dirty="0"/>'
            '<a:t>%s</a:t></a:r></a:p>' % (ppr, int(size * 100),
                                           1 if bold else 0, t))
    return ('<p:sp><p:nvSpPr><p:cNvPr id="%d" name="%s"/>'
            '<p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>'
            '<p:spPr><a:xfrm><a:off x="%d" y="%d"/><a:ext cx="%d" cy="%d"/></a:xfrm>'
            '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
            '<p:txBody><a:bodyPr wrap="square"><a:normAutofit/></a:bodyPr>'
            '<a:lstStyle/>%s</p:txBody></p:sp>'
            % (shape_id, name, x, y, w, h, "".join(paras)))


def _slide_xml(title, lines, has_image=False, image_full=False):
    # 제목: 길이 비례 축소 + 2줄 여유 박스 (화면 밖 이탈 방지)
    title_box = _textbox(2, "title", 457200, 274320, 11277600, 1005840,
                         [(title, _title_size(title), True)])
    if has_image and image_full:
        # 추가 차트 슬라이드: 그림을 크게 중앙 배치 (12:7 비율 유지)
        pic = _picture_xml(4, 1667510, 1417320, 8856980, 5166240)
        return ("""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld %s><p:cSld><p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr/>%s%s</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>"""
                % (_NS, title_box, pic))
    body_paras = []
    base_size = 12 if has_image else 13
    for i, line in enumerate(lines):
        s = str(line)
        is_head = s.startswith("[")
        is_meta = i == 0 and s.startswith("·")
        if is_meta:
            body_paras.append((s, 9, False))          # 메타줄은 작은 회색톤 느낌
        elif is_head:
            body_paras.append((s, base_size + 2, True, 8))  # 섹션 머리글: 위 여백
        else:
            body_paras.append((s, base_size, False))
    if has_image:
        # 차트(좌 6.9") + 인사이트 텍스트(우 5.6") — 세로 정렬 맞춤
        pic = _picture_xml(4, 365760, 1417320, 6309360, 3680460)
        body_box = _textbox(3, "body", 6858000, 1417320, 5029200, 5029200,
                            body_paras or [(" ", base_size, False)])
        shapes = pic + body_box
    else:
        body_box = _textbox(3, "body", 548640, 1417320, 11094720, 5029200,
                            body_paras or [(" ", base_size, False)])
        shapes = body_box
    return ("""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld %s><p:cSld><p:spTree>
<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
<p:grpSpPr/>%s%s</p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:sld>"""
            % (_NS, title_box, shapes))


def _minimal_pptx(slides):
    """외부 의존성 없는 PPTX 생성 (16:9, 텍스트 전용)."""
    n = len(slides)
    ct_overrides = "".join(
        '<Override PartName="/ppt/slides/slide%d.xml" ContentType='
        '"application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        % (i + 1) for i in range(n))
    pres_rels = ['<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/'
                 'officeDocument/2006/relationships/slideMaster" '
                 'Target="slideMasters/slideMaster1.xml"/>']
    sld_ids = []
    for i in range(n):
        rid = "rId%d" % (i + 2)
        pres_rels.append('<Relationship Id="%s" Type="http://schemas.openxmlformats.'
                         'org/officeDocument/2006/relationships/slide" '
                         'Target="slides/slide%d.xml"/>' % (rid, i + 1))
        sld_ids.append('<p:sldId id="%d" r:id="%s"/>' % (256 + i, rid))
    presentation = ("""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation %s><p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
<p:sldIdLst>%s</p:sldIdLst>
<p:sldSz cx="12192000" cy="6858000"/><p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>""" % (_NS, "".join(sld_ids)))
    pres_rels_xml = ('<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                     '<Relationships xmlns="http://schemas.openxmlformats.org/'
                     'package/2006/relationships">%s</Relationships>'
                     % "".join(pres_rels))
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", _CT % ct_overrides)
        z.writestr("_rels/.rels", _ROOT_RELS)
        z.writestr("ppt/presentation.xml", presentation)
        z.writestr("ppt/_rels/presentation.xml.rels", pres_rels_xml)
        z.writestr("ppt/slideMasters/slideMaster1.xml", _MASTER)
        z.writestr("ppt/slideMasters/_rels/slideMaster1.xml.rels", _MASTER_RELS)
        z.writestr("ppt/slideLayouts/slideLayout1.xml", _LAYOUT)
        z.writestr("ppt/slideLayouts/_rels/slideLayout1.xml.rels", _LAYOUT_RELS)
        z.writestr("ppt/theme/theme1.xml", _THEME)
        for i, sl in enumerate(slides):
            has_img = bool(sl.get("image"))
            img_rel = ""
            if has_img:
                media_name = "image%d.%s" % (i + 1, sl.get("ext") or "png")
                z.writestr("ppt/media/%s" % media_name, sl["image"])
                img_rel = _IMG_REL % media_name
            z.writestr("ppt/slides/slide%d.xml" % (i + 1),
                       _slide_xml(sl["title"], sl["lines"], has_image=has_img,
                                  image_full=bool(sl.get("image_full"))))
            z.writestr("ppt/slides/_rels/slide%d.xml.rels" % (i + 1),
                       _SLIDE_RELS % img_rel)
    return buf.getvalue()
